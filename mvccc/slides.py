#!/usr/bin/env python3

# vim: set fileencoding=utf-8 :

from datetime import date, timedelta
from pathlib import Path
from pprint import pformat
from typing import Dict, Generator, List, Tuple

import attr
from absl import app, flags, logging as log
from pptx import Presentation

from bible.index import parse_citations
from bible.scripture import BibleVerse, scripture

flags.DEFINE_bool("extract_only", False, "extract text from pptx")
flags.DEFINE_string("pptx", "", "The pptx")
flags.DEFINE_string("master_pptx", "mvccc_master.pptx", "The template pptx")

flags.DEFINE_string("choir", "", "The hymn by choir")
flags.DEFINE_multi_string("hymns", [], "The hymns by congregation")
flags.DEFINE_string("response", "", "The hymn after the teaching")
flags.DEFINE_string("offering", "", "The hymn for the offering")

flags.DEFINE_string("scripture", "", "The bible scriptures")
flags.DEFINE_string("memorize", "", "The bible scripture to memorize")  # verse of the week

flags.DEFINE_string("message", "", "The message")
flags.DEFINE_string("messager", "", "The messager")

flags.DEFINE_bool("communion", None, "Whether to have communion")

FLAGS = flags.FLAGS

PROCESSED = "processed"

# ------------------------------------------------------------------------------


def next_sunday(today: date = None) -> str:
    if today is None:
        today = date.today()
    sunday = today + timedelta(6 - today.weekday())
    return sunday.isoformat()


def extract_slides_text(ppt: Presentation) -> Generator[Tuple[int, List[List[str]]], None, None]:
    for idx, slide in enumerate(ppt.slides):
        shape_text_list: List[List[str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paragraph_text_list: List[str] = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text_list.append("".join(run.text.replace("\xa0", " ") for run in paragraph.runs))
            while not paragraph_text_list[-1]:
                paragraph_text_list.pop()
            shape_text_list.append(paragraph_text_list)

        yield idx, shape_text_list


# ------------------------------------------------------------------------------

LAYOUT_PRELUDE = 0
LAYOUT_MESSAGE = 1
LAYOUT_HYMN = 2
LAYOUT_SCRIPTURE = 3
LAYOUT_MEMORIZE = 4
LAYOUT_TEACHING = 5
LAYOUT_SECTION = 6
LAYOUT_BLANK = 7


@attr.s
class Prelude:
    message: str = attr.ib()
    picture: str = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_PRELUDE])
        message, picture = slide.placeholders
        message.text = padding + self.message
        picture.insert_picture(self.picture)

        return ppt


@attr.s
class Message:
    message: str = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_MESSAGE])
        message, = slide.placeholders
        message.text = padding + self.message

        return ppt


@attr.s
class Hymn:
    filename: str = attr.ib()  # can be index number, hymn's title
    lyrics: List[Tuple[str, List[str]]] = attr.ib()  # List[title, paragraph]

    def add_to(self, ppt: Presentation, padding: str = " ") -> Presentation:
        for _, (title, paragraph) in self.lyrics:
            slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_HYMN])
            title_holder, paragraph_holder = slide.placeholders
            title_holder.text = title[0]
            # XXX: workaround alignment problem
            paragraph[0] = padding + paragraph[0]
            paragraph_holder.text = "\n".join(paragraph)

        return ppt


def search_hymn_ppt(keyword: str, basepath: Path = None) -> List[Hymn]:
    if basepath is None:
        basepath = Path(PROCESSED)

    keyword = keyword.replace(".pptx", "")
    ptn = f"**/*{keyword}*.pptx"
    glob = basepath.glob(ptn)
    found = list(glob)

    if not found:
        # interchangeability characters
        interchangebles = [("你", "祢", "袮"), ("寶", "寳"), ("他", "祂"), ("于", "於")]
        for t in interchangebles:
            for w in t:
                if w not in ptn:
                    continue
                for w1 in t:
                    if w == w1:
                        continue
                    glob = basepath.glob(ptn.replace(w, w1))
                    found = list(glob)
                    if found:
                        break

    assert found, f"can not find anything match {ptn}."
    if len(found) > 1:
        log.warn(f"found more than 1 files for {ptn}. {[p.as_posix() for p in found]}")

    found = [path for path in found if path.stem == keyword] + [path for path in found if path.stem != keyword]

    result: List[Hymn] = []
    for path in found:
        ppt = Presentation(path.as_posix())
        lyrics = list(extract_slides_text(ppt))
        hymn = Hymn(path.name, lyrics)
        log.info(f"keyword={keyword}, lyrics=\n{pformat(hymn.lyrics)}")
        result.append(hymn)

    return result


@attr.s
class Section:
    title: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_SECTION])
        title, = slide.placeholders
        title.text = self.title

        return ppt


@attr.s
class Scripture:
    citations: str = attr.ib()
    cite_verses: Dict[str, List[BibleVerse]] = attr.ib()

    def add_to(self, ppt: Presentation, padding="  ") -> Presentation:
        for cite, verses in self.cite_verses.items():
            for idx, bv in enumerate(verses):
                if idx % 2 == 0:
                    slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_SCRIPTURE])
                title, message = slide.placeholders
                title.text = cite
                message.text += (padding if idx % 2 == 0 else "\n") + f"{bv.verse}\u3000{bv.text}"

        return ppt


def to_scripture(citations: str) -> Scripture:
    bible = scripture()
    cite_verses = bible.search(parse_citations(citations).items())
    for cite, verses in cite_verses.items():
        log.info(f"citation={cite}, verses=\n{pformat(verses)}")

    return Scripture(citations, cite_verses)


@attr.s
class Memorize:
    citation: str = attr.ib()
    verses: List[BibleVerse] = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_MEMORIZE])
        title, message = slide.placeholders
        title.text = "本週金句"
        message.text = padding + "".join(bv.text for bv in self.verses) + f"\n\n{self.citation:>35}"

        return ppt


@attr.s
class Teaching:
    title: str = attr.ib()
    message: str = attr.ib()
    messenger: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_TEACHING])
        message, = slide.placeholders
        message.text = "\n\n".join([self.title, self.message, self.messenger])

        return ppt


@attr.s
class Blank:
    def add_to(self, ppt: Presentation) -> Presentation:
        _ = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_BLANK])

        return ppt


def mvccc_slides(
    hymns: List[str],
    scripture: str,
    memorize: str,
    message: str,
    messager: str,
    choir: str,
    response: str,
    offering: str,
    communion: bool,
) -> List:
    slides = [
        Prelude("請儘量往前或往中間坐,並將手機關閉或關至靜音,預備心敬拜！", "silence_phone1.png"),
        Message(
            """惟耶和華在他的聖殿中；全地的人，都當在他面前肅敬靜默。

                    哈巴谷書 2:20"""
        ),
    ]
    hymn = search_hymn_ppt("聖哉聖哉聖哉")
    slides.append(hymn[0])

    slides.append(Section("宣  召"))

    slides.append(Section("頌  讚"))
    for kw in hymns:
        r = search_hymn_ppt(kw)
        slides.append(r[0])

    slides.append(Section("祈  禱"))

    slides.append(Section("讀  經"))

    slides.append(to_scripture(scripture))
    for cite, verses in to_scripture(memorize).cite_verses.items():
        slides.append(Memorize(cite, verses))
        break
    slides.append(Blank())

    slides.append(Section("獻  詩"))
    if choir:
        hymn = search_hymn_ppt(choir)[0]
        slides.append(hymn)

    slides.append(Teaching("信息", f"「{message}」", f"{messager}"))

    slides.append(Section("回  應"))
    if response:
        hymn = search_hymn_ppt(response)[0]
        slides.append(hymn)

    if offering:
        hymn = search_hymn_ppt(offering)[0]
        slides.append(hymn)

    slides.append(Section("奉 獻 禱 告"))

    if communion:
        slides.append(Section("聖  餐"))

    slides.append(Section("歡 迎 您"))
    slides.append(Section("家 事 分 享"))

    hymn = search_hymn_ppt("三一頌")[0]
    slides.append(hymn)

    slides.append(Section("祝  福"))
    slides.append(Section("默  禱"))
    slides.append(Blank())

    return slides


def to_pptx(slides: List, master_slide: Presentation) -> Presentation:
    ppt = master_slide

    for slide in slides:
        slide.add_to(ppt)

    return ppt


# ------------------------------------------------------------------------------


def main(argv):
    del argv

    if FLAGS.extract_only:
        ppt = Presentation(FLAGS.pptx)

        for idx, text in extract_slides_text(ppt):
            print(f"{idx+1:02d} {text}")
        return

    slides = mvccc_slides(
        hymns=FLAGS.hymns,
        scripture=FLAGS.scripture,
        memorize=FLAGS.memorize,
        message=FLAGS.message,
        messager=FLAGS.messager,
        choir=FLAGS.choir,
        response=FLAGS.response,
        offering=FLAGS.offering,
        communion=FLAGS.communion,
    )
    master = Presentation(FLAGS.master_pptx)
    ppt = to_pptx(slides, master)
    ppt.save(FLAGS.pptx)


if __name__ == "__main__":
    app.run(main)
